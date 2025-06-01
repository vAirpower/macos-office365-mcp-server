"""
PowerPoint Controller for Office 365 MCP Server
Handles all PowerPoint automation operations on macOS.
"""

import asyncio
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from integrations.applescript_bridge import AppleScriptBridge
from utils.logger import setup_logger

logger = setup_logger(__name__)

class PowerPointController:
    """Controller for PowerPoint operations using both AppleScript and python-pptx."""
    
    def __init__(self):
        self.applescript = AppleScriptBridge()
        self.active_presentations: Dict[str, Dict[str, Any]] = {}
        self.temp_dir = Path.home() / "tmp" / "office365_mcp"
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    async def create_presentation(
        self,
        title: str,
        theme: str = "default",
        template_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """Create a new PowerPoint presentation.
        
        Args:
            title: Presentation title
            theme: Theme name or path
            template_path: Optional template file path
            
        Returns:
            Dict with presentation metadata
        """
        try:
            presentation_id = str(uuid.uuid4())
            
            # Create presentation using python-pptx
            if template_path and Path(template_path).exists():
                prs = Presentation(template_path)
            else:
                prs = Presentation()
            
            # Set title if first slide exists
            if prs.slides:
                title_slide = prs.slides[0]
                if title_slide.shapes.title:
                    title_slide.shapes.title.text = title
            else:
                # Add title slide
                title_slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(title_slide_layout)
                slide.shapes.title.text = title
            
            # Save temporary file
            temp_file = self.temp_dir / f"{presentation_id}.pptx"
            prs.save(str(temp_file))
            
            # Try to open in PowerPoint via AppleScript
            applescript_success = False
            try:
                await self.applescript.open_powerpoint_file(str(temp_file))
                applescript_success = True
            except Exception as e:
                logger.warning(f"Could not open in PowerPoint app: {e}")
            
            # Store presentation metadata
            presentation_data = {
                "presentation_id": presentation_id,
                "title": title,
                "theme": theme,
                "file_path": str(temp_file),
                "slide_count": len(prs.slides),
                "applescript_available": applescript_success,
                "created_at": asyncio.get_event_loop().time()
            }
            
            self.active_presentations[presentation_id] = {
                "metadata": presentation_data,
                "pptx_object": prs,
                "slides": {}
            }
            
            logger.info(f"Created presentation: {title} ({presentation_id})")
            return presentation_data
            
        except Exception as e:
            logger.error(f"Failed to create presentation: {e}")
            raise
    
    async def add_slide(
        self,
        presentation_id: str,
        layout: str = "Title and Content",
        position: Optional[int] = None
    ) -> Dict[str, Any]:
        """Add a new slide to a presentation.
        
        Args:
            presentation_id: ID of the presentation
            layout: Slide layout name
            position: Position to insert slide
            
        Returns:
            Dict with slide metadata
        """
        try:
            if presentation_id not in self.active_presentations:
                raise ValueError(f"Presentation {presentation_id} not found")
            
            prs = self.active_presentations[presentation_id]["pptx_object"]
            slide_id = str(uuid.uuid4())
            
            # Map layout names to indices
            layout_map = {
                "Title Slide": 0,
                "Title and Content": 1,
                "Section Header": 2,
                "Two Content": 3,
                "Comparison": 4,
                "Title Only": 5,
                "Blank": 6,
                "Content with Caption": 7,
                "Picture with Caption": 8
            }
            
            layout_index = layout_map.get(layout, 1)  # Default to Title and Content
            slide_layout = prs.slide_layouts[layout_index]
            
            # Add slide
            if position is not None and 0 <= position <= len(prs.slides):
                # Insert at specific position (requires manual reordering)
                slide = prs.slides.add_slide(slide_layout)
                # Note: python-pptx doesn't support direct insertion at position
                # This would require AppleScript for precise positioning
            else:
                slide = prs.slides.add_slide(slide_layout)
            
            slide_index = len(prs.slides) - 1
            
            # Store slide metadata
            slide_data = {
                "slide_id": slide_id,
                "presentation_id": presentation_id,
                "layout": layout,
                "index": slide_index,
                "created_at": asyncio.get_event_loop().time()
            }
            
            self.active_presentations[presentation_id]["slides"][slide_id] = {
                "metadata": slide_data,
                "slide_object": slide
            }
            
            # Update presentation metadata
            self.active_presentations[presentation_id]["metadata"]["slide_count"] = len(prs.slides)
            
            # Save updated presentation
            temp_file = self.active_presentations[presentation_id]["metadata"]["file_path"]
            prs.save(temp_file)
            
            logger.info(f"Added slide to presentation {presentation_id}")
            return slide_data
            
        except Exception as e:
            logger.error(f"Failed to add slide: {e}")
            raise
    
    async def add_text(
        self,
        slide_id: str,
        text: str,
        placeholder: str = "content",
        formatting: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """Add text to a slide.
        
        Args:
            slide_id: ID of the slide
            text: Text content
            placeholder: Placeholder type (title, content, etc.)
            formatting: Text formatting options
            
        Returns:
            Dict with operation status
        """
        try:
            # Find slide
            slide_obj = None
            presentation_id = None
            
            for pres_id, pres_data in self.active_presentations.items():
                if slide_id in pres_data["slides"]:
                    slide_obj = pres_data["slides"][slide_id]["slide_object"]
                    presentation_id = pres_id
                    break
            
            if not slide_obj:
                raise ValueError(f"Slide {slide_id} not found")
            
            # Add text based on placeholder type
            if placeholder == "title" and slide_obj.shapes.title:
                text_frame = slide_obj.shapes.title.text_frame
                text_frame.text = text
            elif placeholder == "content":
                # Find content placeholder or add text box
                content_placeholder = None
                for shape in slide_obj.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    text_frame = content_placeholder.text_frame
                    text_frame.text = text
                else:
                    # Add text box
                    left = Inches(1)
                    top = Inches(1.5)
                    width = Inches(8)
                    height = Inches(5)
                    textbox = slide_obj.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text
            else:
                # Add as text box
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(5)
                textbox = slide_obj.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text
            
            # Apply formatting if provided
            if formatting and text_frame:
                await self._apply_text_formatting(text_frame, formatting)
            
            # Save presentation
            prs = self.active_presentations[presentation_id]["pptx_object"]
            temp_file = self.active_presentations[presentation_id]["metadata"]["file_path"]
            prs.save(temp_file)
            
            logger.info(f"Added text to slide {slide_id}")
            return {
                "status": "success",
                "slide_id": slide_id,
                "text_length": len(text),
                "placeholder": placeholder
            }
            
        except Exception as e:
            logger.error(f"Failed to add text: {e}")
            raise
    
    async def add_image(
        self,
        slide_id: str,
        image_source: str,
        position: Dict[str, float],
        size: Dict[str, float]
    ) -> Dict[str, Any]:
        """Add an image to a slide.
        
        Args:
            slide_id: ID of the slide
            image_source: Path to image file or URL
            position: Position dict with x, y coordinates
            size: Size dict with width, height
            
        Returns:
            Dict with operation status
        """
        try:
            # Find slide
            slide_obj = None
            presentation_id = None
            
            for pres_id, pres_data in self.active_presentations.items():
                if slide_id in pres_data["slides"]:
                    slide_obj = pres_data["slides"][slide_id]["slide_object"]
                    presentation_id = pres_id
                    break
            
            if not slide_obj:
                raise ValueError(f"Slide {slide_id} not found")
            
            # Handle image source
            if image_source.startswith(("http://", "https://")):
                # Download image (simplified - would need proper download logic)
                raise NotImplementedError("URL image download not yet implemented")
            else:
                # Local file
                image_path = Path(image_source)
                if not image_path.exists():
                    raise FileNotFoundError(f"Image file not found: {image_source}")
            
            # Add image to slide
            left = Inches(position.get("x", 1))
            top = Inches(position.get("y", 1))
            width = Inches(size.get("width", 4))
            height = Inches(size.get("height", 3))
            
            slide_obj.shapes.add_picture(str(image_path), left, top, width, height)
            
            # Save presentation
            prs = self.active_presentations[presentation_id]["pptx_object"]
            temp_file = self.active_presentations[presentation_id]["metadata"]["file_path"]
            prs.save(temp_file)
            
            logger.info(f"Added image to slide {slide_id}")
            return {
                "status": "success",
                "slide_id": slide_id,
                "image_source": image_source,
                "position": position,
                "size": size
            }
            
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            raise
    
    async def add_speaker_notes(
        self,
        slide_id: str,
        notes: str
    ) -> Dict[str, Any]:
        """Add speaker notes to a slide.
        
        Args:
            slide_id: ID of the slide
            notes: Speaker notes content
            
        Returns:
            Dict with operation status
        """
        try:
            # Find slide
            slide_obj = None
            presentation_id = None
            
            for pres_id, pres_data in self.active_presentations.items():
                if slide_id in pres_data["slides"]:
                    slide_obj = pres_data["slides"][slide_id]["slide_object"]
                    presentation_id = pres_id
                    break
            
            if not slide_obj:
                raise ValueError(f"Slide {slide_id} not found")
            
            # Add notes
            notes_slide = slide_obj.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            
            # Save presentation
            prs = self.active_presentations[presentation_id]["pptx_object"]
            temp_file = self.active_presentations[presentation_id]["metadata"]["file_path"]
            prs.save(temp_file)
            
            logger.info(f"Added speaker notes to slide {slide_id}")
            return {
                "status": "success",
                "slide_id": slide_id,
                "notes_length": len(notes)
            }
            
        except Exception as e:
            logger.error(f"Failed to add speaker notes: {e}")
            raise
    
    async def save_presentation(
        self,
        presentation_id: str,
        file_path: str,
        format: str = "pptx"
    ) -> Dict[str, Any]:
        """Save a presentation to file.
        
        Args:
            presentation_id: ID of the presentation
            file_path: Path to save the file
            format: File format (pptx, pdf, etc.)
            
        Returns:
            Dict with operation status
        """
        try:
            if presentation_id not in self.active_presentations:
                raise ValueError(f"Presentation {presentation_id} not found")
            
            prs = self.active_presentations[presentation_id]["pptx_object"]
            save_path = Path(file_path)
            
            # Ensure directory exists
            save_path.parent.mkdir(parents=True, exist_ok=True)
            
            if format.lower() == "pptx":
                prs.save(str(save_path))
            elif format.lower() == "pdf":
                # PDF export would require AppleScript or additional libraries
                # For now, save as PPTX and note the limitation
                pptx_path = save_path.with_suffix(".pptx")
                prs.save(str(pptx_path))
                logger.warning(f"PDF export not implemented, saved as PPTX: {pptx_path}")
                save_path = pptx_path
            else:
                raise ValueError(f"Unsupported format: {format}")
            
            # Update metadata
            self.active_presentations[presentation_id]["metadata"]["file_path"] = str(save_path)
            
            logger.info(f"Saved presentation to {save_path}")
            return {
                "status": "success",
                "presentation_id": presentation_id,
                "file_path": str(save_path),
                "format": format
            }
            
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise
    
    async def _apply_text_formatting(
        self,
        text_frame,
        formatting: Dict[str, Any]
    ) -> None:
        """Apply text formatting to a text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            formatting: Formatting options
        """
        try:
            paragraph = text_frame.paragraphs[0]
            
            # Font size
            if "font_size" in formatting:
                paragraph.font.size = Pt(formatting["font_size"])
            
            # Font name
            if "font_name" in formatting:
                paragraph.font.name = formatting["font_name"]
            
            # Bold
            if "bold" in formatting:
                paragraph.font.bold = formatting["bold"]
            
            # Italic
            if "italic" in formatting:
                paragraph.font.italic = formatting["italic"]
            
            # Color
            if "color" in formatting:
                color = formatting["color"]
                if isinstance(color, str) and color.startswith("#"):
                    # Convert hex to RGB
                    hex_color = color.lstrip("#")
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    paragraph.font.color.rgb = RGBColor(*rgb)
            
            # Alignment
            if "alignment" in formatting:
                align_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY
                }
                if formatting["alignment"] in align_map:
                    paragraph.alignment = align_map[formatting["alignment"]]
            
        except Exception as e:
            logger.warning(f"Failed to apply some text formatting: {e}")
    
    async def get_presentation_info(self, presentation_id: str) -> Dict[str, Any]:
        """Get information about a presentation.
        
        Args:
            presentation_id: ID of the presentation
            
        Returns:
            Dict with presentation information
        """
        if presentation_id not in self.active_presentations:
            raise ValueError(f"Presentation {presentation_id} not found")
        
        return self.active_presentations[presentation_id]["metadata"]
    
    async def list_presentations(self) -> List[Dict[str, Any]]:
        """List all active presentations.
        
        Returns:
            List of presentation metadata
        """
        return [data["metadata"] for data in self.active_presentations.values()]
